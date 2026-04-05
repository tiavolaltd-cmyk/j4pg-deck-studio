"""
J4PG — Fill Engine
Remplit le master template TOP20/TOP12 avec les données d'un joueur.
Usage :
    from fill_engine import fill_deck
    fill_deck("master_top20.pptx", data_dict, images_dict, "output.pptx")
"""

import copy, io, re, os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor

TAG_RE = re.compile(r'\[\[([A-Z0-9_]+)\]\]')

# ─── Constantes barre ──────────────────────────────────────────────────────────
S3_BAR_MAX_WIDTH_IN = 1.253   # 100% atteint = 1.253"
S4_FORCE_BAR_MAX_IN = 3.800   # Forces (vert) — max width
S4_AXE_BAR_MAX_IN   = 3.900   # Axes d'amélioration (rouge) — max width


# ═══════════════════════════════════════════════════════════════════════════════
# 1. REMPLACEMENT TEXTE — parcours tous les text frames
# ═══════════════════════════════════════════════════════════════════════════════

def _replace_tags_in_run(run, data: dict) -> None:
    """Remplace les [[TAG]] dans un run en préservant le formatage."""
    text = run.text
    def replacer(m):
        tag = m.group(1)
        return str(data.get(tag, f"[[{tag}]]"))
    run.text = TAG_RE.sub(replacer, text)


def replace_all_text_tags(prs: Presentation, data: dict) -> None:
    """Parcourt toutes les slides et remplace les [[TAG]] dans les runs."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if TAG_RE.search(run.text):
                        _replace_tags_in_run(run, data)
            # Certains shapes ont aussi des placeholder dans le titre
            if hasattr(shape, "text") and "[[" in shape.text:
                # Fallback: forcer le remplacement dans tous les runs
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        _replace_tags_in_run(run, data)


# ═══════════════════════════════════════════════════════════════════════════════
# 2. REDIMENSIONNEMENT BARRES
# ═══════════════════════════════════════════════════════════════════════════════

def _emu(inches: float) -> int:
    return int(inches * 914400)


def resize_s3_bars(prs: Presentation, kpi_vals: list) -> None:
    """
    kpi_vals : liste de 8 tuples (player_val, target_val) — float/float.
    La barre s3_kpi{i}_bar est redimensionnée proportionnellement.
    Si player_val >= target_val → largeur max (100%).
    """
    slide = prs.slides[2]
    for i, (player_val, target_val) in enumerate(kpi_vals, 1):
        bar_name = f"s3_kpi{i}_bar"
        for shape in slide.shapes:
            if shape.name == bar_name:
                if target_val and target_val > 0:
                    ratio = min(float(player_val) / float(target_val), 1.0)
                else:
                    ratio = 0.0
                new_width = _emu(ratio * S3_BAR_MAX_WIDTH_IN)
                shape.width = max(new_width, _emu(0.02))  # min 0.02" pour rester visible
                break


def resize_s4_bars(prs: Presentation,
                   force_percentiles: list,
                   axe_percentiles: list) -> None:
    """
    force_percentiles : liste de 3 float (0–100) pour les forces (barres vertes).
    axe_percentiles   : liste de 3 float (0–100) pour les axes d'amélioration.
    """
    slide = prs.slides[3]
    for i, pct in enumerate(force_percentiles, 1):
        for shape in slide.shapes:
            if shape.name == f"s4_f{i}_bar":
                ratio = max(0.0, min(float(pct) / 100.0, 1.0))
                shape.width = max(_emu(ratio * S4_FORCE_BAR_MAX_IN), _emu(0.02))
                break
    for i, pct in enumerate(axe_percentiles, 1):
        for shape in slide.shapes:
            if shape.name == f"s4_a{i}_bar":
                ratio = max(0.0, min(float(pct) / 100.0, 1.0))
                shape.width = max(_emu(ratio * S4_AXE_BAR_MAX_IN), _emu(0.02))
                break


# ═══════════════════════════════════════════════════════════════════════════════
# 3. INSERTION IMAGES
# ═══════════════════════════════════════════════════════════════════════════════

from pptx.util import Pt
from PIL import Image as PILImage


def _get_image_bytes(img_src) -> bytes:
    """Accepte: bytes, file-like object, ou chemin fichier."""
    if isinstance(img_src, (bytes, bytearray)):
        return bytes(img_src)
    if hasattr(img_src, 'read'):
        return img_src.read()
    with open(img_src, 'rb') as f:
        return f.read()


def _has_blip_fill(shape) -> bool:
    """Retourne True si le shape utilise une image comme fill (blipFill)."""
    for elem in shape._element.iter():
        if 'blipFill' in elem.tag:
            return True
    return False


def _replace_blip_fill(slide, shape, img_bytes: bytes) -> bool:
    """
    Remplace l'image dans un shape blipFill (ex: PHOTO_PLAYER = Freeform masque).
    Met a jour l'attribut r:embed pour pointer vers la nouvelle image.
    """
    NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    REL_TYPE = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image'
    img_io = io.BytesIO(img_bytes)
    image_part = slide.part.package.get_or_add_image_part(img_io)
    new_rId = slide.part.relate_to(image_part, REL_TYPE)
    embed_attr = f'{{{NS_R}}}embed'
    for elem in shape._element.iter():
        if elem.get(embed_attr) is not None:
            elem.set(embed_attr, new_rId)
            return True
    return False


def _insert_image_at_shape(slide, shape_name: str, img_bytes: bytes,
                            keep_aspect: bool = False) -> bool:
    """
    Remplace le shape nomme `shape_name` par une image.
    Cas 1 : shape blipFill (ex: PHOTO_PLAYER) -> swap r:embed, conserve la forme.
    Cas 2 : shape Picture ou rectangle -> suppression + ajout Picture classique.
    Si keep_aspect=True : l'image est centree dans le shape sans deformation.
    Retourne True si le shape a ete trouve.
    """
    from pptx.util import Inches
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    target = None
    for shape in slide.shapes:
        if shape.name == shape_name:
            target = shape
            break
    if target is None:
        return False

    # Cas 1 : shape avec blipFill (Freeform masque comme PHOTO_PLAYER)
    if _has_blip_fill(target):
        return _replace_blip_fill(slide, target, img_bytes)

    left, top, width, height = target.left, target.top, target.width, target.height

    if keep_aspect:
        # Calcul des dimensions réelles de l'image
        img_io = io.BytesIO(img_bytes)
        with PILImage.open(img_io) as pil_img:
            iw, ih = pil_img.size
        ratio_img   = iw / ih
        ratio_box   = width / height
        if ratio_img > ratio_box:
            # Image plus large → ajuster la hauteur
            new_w = width
            new_h = int(width / ratio_img)
            top  += (height - new_h) // 2
            height = new_h
        else:
            # Image plus haute → ajuster la largeur
            new_h = height
            new_w = int(height * ratio_img)
            left += (width - new_w) // 2
            width = new_w

    # Supprimer l'ancien shape
    sp = target._element
    sp.getparent().remove(sp)

    # Ajouter la nouvelle image
    img_io = io.BytesIO(img_bytes)
    pic = slide.shapes.add_picture(img_io, left, top, width, height)
    pic.name = shape_name
    return True


def insert_images(prs: Presentation, images: dict) -> None:
    """
    images : dict { shape_name: img_src }
    Règles spéciales :
      - IMG_HEATMAP → keep_aspect=True (pas de déformation)
      - Autres       → étirer au shape (comportement PowerPoint standard)
    """
    HEATMAP_KEEP_ASPECT = {"IMG_HEATMAP"}

    slide_map = {}  # shape_name → index slide
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.name in images:
                slide_map[shape.name] = i

    for shape_name, img_src in images.items():
        if img_src is None:
            continue
        slide_idx = slide_map.get(shape_name)
        if slide_idx is None:
            continue
        slide = prs.slides[slide_idx]
        img_bytes = _get_image_bytes(img_src)
        keep = shape_name in HEATMAP_KEEP_ASPECT
        ok = _insert_image_at_shape(slide, shape_name, img_bytes, keep_aspect=keep)
        if not ok:
            print(f"  [WARN] Shape non trouvé : {shape_name}")


# ═══════════════════════════════════════════════════════════════════════════════
# 4. POINT D'ENTRÉE PRINCIPAL
# ═══════════════════════════════════════════════════════════════════════════════

def fill_deck(template_path: str,
              data: dict,
              images: dict,
              output_path: str,
              s3_kpi_vals: list = None,
              s4_force_pcts: list = None,
              s4_axe_pcts: list = None) -> str:
    """
    Remplit le master template et sauvegarde le deck final.

    Paramètres
    ----------
    template_path  : chemin vers master_top20.pptx (ou top12)
    data           : dict { 'TAG': 'valeur' } — toutes les valeurs texte
    images         : dict { 'SHAPE_NAME': img_bytes_ou_chemin }
    output_path    : chemin de sortie
    s3_kpi_vals    : liste de 8 tuples (player_val, target_val)
    s4_force_pcts  : liste de 3 percentiles (0–100) pour les forces
    s4_axe_pcts    : liste de 3 percentiles (0–100) pour les axes
    """
    prs = Presentation(template_path)

    # 1. Textes
    replace_all_text_tags(prs, data)

    # 2. Barres slide 3
    if s3_kpi_vals and len(s3_kpi_vals) == 8:
        resize_s3_bars(prs, s3_kpi_vals)

    # 3. Barres slide 4
    if s4_force_pcts and s4_axe_pcts:
        resize_s4_bars(prs, s4_force_pcts, s4_axe_pcts)

    # 4. Images
    if images:
        insert_images(prs, images)

    prs.save(output_path)
    return output_path


# ═══════════════════════════════════════════════════════════════════════════════
# 5. TEST RAPIDE — Nathan Bitumazala (MOA)
# ═══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    import os

    TEMPLATE = "/sessions/adoring-hopeful-euler/deck_work/master_top20.pptx"
    OUTPUT   = "/sessions/adoring-hopeful-euler/deck_work/test_nathan_filled.pptx"

    # Données de test Nathan Bitumazala
    DATA = {
        # Slide 1 — Cover
        "PRENOM_UPPER": "NATHAN",
        "NOM_UPPER":    "BITUMAZALA",
        "POSTE":        "Milieu Offensif / Ailier",
        "AGE":          "22",
        "TAILLE":       "178",
        "PIED":         "Droit",
        "CLUB":         "AS Nancy-Lorraine",
        "COMPETITION":  "Ligue 2",
        "DECK_TYPE":    "ZOOM JOUEUR — TOP 20",
        "SAISON":       "2024-25",
        "FENETRE":      "J1-J26",

        # Slide 2 — Profil
        "LIGUE":        "Ligue 2",
        "MINUTES":      "1 847",
        "MATCHS":       "22",
        "S2_KPI1_LABEL": "Occasions de but",
        "S2_KPI1_VAL":   "3.42",
        "S2_KPI1_PCTILE":"78e",
        "S2_KPI2_LABEL": "Occasions créées",
        "S2_KPI2_VAL":   "2.87",
        "S2_KPI2_PCTILE":"71e",
        "S2_KPI3_LABEL": "Passes décisives",
        "S2_KPI3_VAL":   "0.43",
        "S2_KPI3_PCTILE":"69e",
        "S2_KPI4_LABEL": "Dribbles réussis (%)",
        "S2_KPI4_VAL":   "54%",
        "S2_KPI4_PCTILE":"65e",
        "S2_KPI5_LABEL": "Tirs cadrés (%)",
        "S2_KPI5_VAL":   "41%",
        "S2_KPI5_PCTILE":"72e",
        "S2_KPI6_LABEL": "Centres réussis (%)",
        "S2_KPI6_VAL":   "28%",
        "S2_KPI6_PCTILE":"58e",
        "S2_KPI7_LABEL": "Passes clés",
        "S2_KPI7_VAL":   "1.91",
        "S2_KPI8_LABEL": "Duels offensifs gagnés (%)",
        "S2_KPI8_VAL":   "47%",
        "S2_KPI9_LABEL": "Passes avant tir",
        "S2_KPI9_VAL":   "2.15",
        "SOURCE":        "SportsBase",

        # Slide 3 — Standard Cible
        "S3_TITLE":      "STANDARD CIBLE",
        "S3_SUBTITLE":   "Milieu Offensif · Ligue 2 · Top 25% référence",
        "S3_KPI1_LABEL": "Occasions de but",
        "S3_KPI1_VAL":   "3.42",
        "S3_KPI1_TARGET":"4.10",
        "S3_KPI1_GAP":   "-0.68",
        "S3_KPI2_LABEL": "Occasions créées",
        "S3_KPI2_VAL":   "2.87",
        "S3_KPI2_TARGET":"3.50",
        "S3_KPI2_GAP":   "-0.63",
        "S3_KPI3_LABEL": "Dribbles réussis (%)",
        "S3_KPI3_VAL":   "54%",
        "S3_KPI3_TARGET":"62%",
        "S3_KPI3_GAP":   "-8%",
        "S3_KPI4_LABEL": "Tirs cadrés (%)",
        "S3_KPI4_VAL":   "41%",
        "S3_KPI4_TARGET":"48%",
        "S3_KPI4_GAP":   "-7%",
        "S3_KPI5_LABEL": "Passes clés",
        "S3_KPI5_VAL":   "1.91",
        "S3_KPI5_TARGET":"1.91",
        "S3_KPI5_GAP":   "0.00",
        "S3_KPI6_LABEL": "Centres réussis (%)",
        "S3_KPI6_VAL":   "28%",
        "S3_KPI6_TARGET":"34%",
        "S3_KPI6_GAP":   "-6%",
        "S3_KPI7_LABEL": "Passes décisives",
        "S3_KPI7_VAL":   "0.43",
        "S3_KPI7_TARGET":"0.43",
        "S3_KPI7_GAP":   "0.00",
        "S3_KPI8_LABEL": "Entrées DT par conduite",
        "S3_KPI8_VAL":   "1.20",
        "S3_KPI8_TARGET":"1.60",
        "S3_KPI8_GAP":   "-0.40",

        # Slide 4 — Analyse Individuelle
        "S4_TITLE":    "ANALYSE INDIVIDUELLE",
        "S4_SUBTITLE": "Forces & Axes d'amélioration",
        "S4_F1_LABEL": "Création offensive",
        "S4_F1_DESC":  "3.42 occ. de but /90 · 2.87 occ. créées /90",
        "S4_F1_PCT":   "78",
        "S4_F1_BADGE": "TOP 25%",
        "S4_F2_LABEL": "Duel balle au pied",
        "S4_F2_DESC":  "54% dribbles réussis · fort au duel offensif",
        "S4_F2_PCT":   "69",
        "S4_F2_BADGE": "TOP 35%",
        "S4_F3_LABEL": "Transition rapide",
        "S4_F3_DESC":  "2.15 passes avant tir /90 · engage vite",
        "S4_F3_PCT":   "60",
        "S4_F3_BADGE": "TOP 40%",
        "S4_A1_LABEL": "Précision tir",
        "S4_A1_DESC":  "41% tirs cadrés · sous la référence top 25%",
        "S4_A1_PCT":   "38",
        "S4_A1_BADGE": "TOP 62%",
        "S4_A2_LABEL": "Centre depuis le couloir",
        "S4_A2_DESC":  "28% centres réussis · volume trop faible",
        "S4_A2_PCT":   "45",
        "S4_A2_BADGE": "TOP 55%",
        "S4_A3_LABEL": "Entrées derniers tiers",
        "S4_A3_DESC":  "1.20 /90 · potentiel exploitation couloir",
        "S4_A3_PCT":   "49",
        "S4_A3_BADGE": "TOP 51%",

        # Slide 5 — Benchmark Ligue
        "S5_TITLE":          "BENCHMARK LIGUE",
        "S5_SUBTITLE":       "Milieu Offensif · Ligue 2 · /90 min",
        "S5_RADAR_SUBTITLE": "Nathan vs. Top 25% Ligue 2",
        "S5_KPI1_LABEL": "Dribbles réussis (%)", "S5_KPI1_N": "54%",  "S5_KPI1_L": "62%",  "S5_KPI1_PCT": "65",
        "S5_KPI2_LABEL": "xA /90",               "S5_KPI2_N": "0.08", "S5_KPI2_L": "0.11", "S5_KPI2_PCT": "61",
        "S5_KPI3_LABEL": "xG /90",               "S5_KPI3_N": "0.22", "S5_KPI3_L": "0.28", "S5_KPI3_PCT": "68",
        "S5_KPI4_LABEL": "Actions surface",       "S5_KPI4_N": "3.10", "S5_KPI4_L": "4.20", "S5_KPI4_PCT": "58",
        "S5_KPI5_LABEL": "Tirs /90",              "S5_KPI5_N": "2.40", "S5_KPI5_L": "2.90", "S5_KPI5_PCT": "72",
        "S5_KPI6_LABEL": "Passes clés /90",       "S5_KPI6_N": "1.91", "S5_KPI6_L": "2.10", "S5_KPI6_PCT": "71",

        # Slide 6 — Zones d'Activité
        "S6_TITLE":    "ZONES D'ACTIVITÉ",
        "S6_SUBTITLE": "Positionnement moyen · J1–J26 · Ligue 2 2024-25",

        # Slide 7 — Comparaison
        "S7_TITLE": "COMPARAISON JOUEURS",
        "J1_NOM": "N. Bitumazala", "J2_NOM": "T. Lemaire", "J2_CLUB": "Le Havre",
        "J3_NOM": "A. Koné", "J3_CLUB": "Valenciennes",
        "J1_NOM_COURT": "Bitumazala", "J1_POSTE": "MOA", "J1_MIN": "1847",
        "J2_NOM_COURT": "Lemaire",   "J2_POSTE": "MOA", "J2_MIN": "1620",
        "J3_NOM_COURT": "Koné",      "J3_POSTE": "MOA", "J3_MIN": "1540",
        "S7_KPI1_LABEL": "Occasions de but",  "S7_KPI1_J1": "3.42", "S7_KPI1_J2": "3.85", "S7_KPI1_J3": "2.91",
        "S7_KPI2_LABEL": "Occasions créées",  "S7_KPI2_J1": "2.87", "S7_KPI2_J2": "3.10", "S7_KPI2_J3": "2.55",
        "S7_KPI3_LABEL": "Passes décisives",  "S7_KPI3_J1": "0.43", "S7_KPI3_J2": "0.51", "S7_KPI3_J3": "0.38",
        "S7_KPI4_LABEL": "Dribbles réussis%", "S7_KPI4_J1": "54%",  "S7_KPI4_J2": "48%",  "S7_KPI4_J3": "51%",
        "S7_KPI5_LABEL": "Tirs cadrés%",      "S7_KPI5_J1": "41%",  "S7_KPI5_J2": "44%",  "S7_KPI5_J3": "39%",
        "S7_KPI6_LABEL": "Centres réussis%",  "S7_KPI6_J1": "28%",  "S7_KPI6_J2": "31%",  "S7_KPI6_J3": "24%",
        "S7_KPI7_LABEL": "Passes avant tir",  "S7_KPI7_J1": "2.15", "S7_KPI7_J2": "2.02", "S7_KPI7_J3": "1.87",
        "S7_KPI8_LABEL": "Entrées DT",        "S7_KPI8_J1": "1.20", "S7_KPI8_J2": "1.45", "S7_KPI8_J3": "1.11",
        "S7_LEGENDE": "N = Nathan Bitumazala · Valeurs /90 min · SportsBase 2024-25",
        "SAISON":      "2024-25",

        # Slide 8 — Axes de Progression
        "S8_TITLE":    "AXES DE PROGRESSION",
        "S8_SUBTITLE": "Plan de développement individuel",
        "S8_P1_TITLE": "AXE 1 · Efficacité devant le but",
        "S8_P1_KPI":   "Tirs cadrés (%)",
        "S8_P1_TODAY": "41%",
        "S8_P1_CAP1":  "46%",
        "S8_P1_CAPOPT":"50%",
        "S8_P1_IMPACT":"+ 2–3 buts/saison",
        "S8_P2_TITLE": "AXE 2 · Exploitation des couloirs",
        "S8_P2_KPI":   "Centres réussis (%)",
        "S8_P2_TODAY": "28%",
        "S8_P2_CAP1":  "32%",
        "S8_P2_CAPOPT":"36%",
        "S8_P2_IMPACT":"+ volume centres exploitables",
        "S8_P3_TITLE": "AXE 3 · Prise de profondeur",
        "S8_P3_KPI":   "Entrées DT par conduite",
        "S8_P3_TODAY": "1.20",
        "S8_P3_CAP1":  "1.40",
        "S8_P3_CAPOPT":"1.65",
        "S8_P3_IMPACT":"Déstabilisation défense adverse",

        # Slide 9 — Plan développement
        "S9_BASE_VALS":   "Tirs cadrés 41% · Centres 28% · Entrées DT 1.20",
        "S9_CAP1_VALS":   "Tirs cadrés 46% · Centres 32% · Entrées DT 1.40",
        "S9_CAP2_VALS":   "Tirs cadrés 48% · Centres 34% · Entrées DT 1.52",
        "S9_CAPOPT_VALS": "Tirs cadrés 50% · Centres 36% · Entrées DT 1.65",
        "S9_IMPACT_BUTS": "+2/3 buts",
        "S9_IMPACT_PASSES":"+ passes décisives",
    }

    # Barres slide 3 : (val_joueur, val_cible)
    S3_VALS = [
        (3.42, 4.10),   # Occasions de but
        (2.87, 3.50),   # Occasions créées
        (54,   62),     # Dribbles %
        (41,   48),     # Tirs cadrés %
        (1.91, 1.91),   # Passes clés — 100%
        (28,   34),     # Centres %
        (0.43, 0.43),   # Passes décisives — 100%
        (1.20, 1.60),   # Entrées DT
    ]

    # Barres slide 4 : percentiles
    FORCE_PCTS = [78, 69, 60]
    AXE_PCTS   = [38, 45, 49]

    out = fill_deck(
        template_path=TEMPLATE,
        data=DATA,
        images={},  # Pas d'images en test
        output_path=OUTPUT,
        s3_kpi_vals=S3_VALS,
        s4_force_pcts=FORCE_PCTS,
        s4_axe_pcts=AXE_PCTS,
    )

    print(f"✓ Deck test généré : {out}")
    print(f"  Taille : {os.path.getsize(out):,} bytes")

    # Vérification : s'assurer qu'il ne reste aucun [[TAG]] non remplacé
    prs_check = Presentation(out)
    remaining = []
    for i, slide in enumerate(prs_check.slides, 1):
        for shape in slide.shapes:
            if hasattr(shape, 'text') and '[[' in shape.text:
                remaining.append(f"S{i}/{shape.name}: {shape.text[:60]}")
    if remaining:
        print(f"\n  ⚠ Tags non remplacés ({len(remaining)}) :")
        for r in remaining[:10]:
            print(f"    {r}")
    else:
        print("  ✓ Aucun tag résiduel")
